import os
from pptx import Presentation
import logging
from typing import List
from .data_room_schemas import ExtractedMetric

logger = logging.getLogger(__name__)

def parse_pitch_deck(file_path: str, document_name: str) -> List[ExtractedMetric]:
    """
    Parses a PPTX pitch deck, extracts text from slides and speaker notes,
    and uses the LLM router to extract key metrics or claims.
    """
    if not os.path.exists(file_path):
        logger.warning(f"File not found: {file_path}")
        return []

    try:
        if not file_path.endswith(".pptx"):
            logger.warning(f"Unsupported deck format: {file_path}")
            return []

        prs = Presentation(file_path)
        
        slide_texts = []
        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
            
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                
            slide_content = f"Slide {i+1}: " + " ".join(text_runs)
            if notes:
                slide_content += f" | Speaker Notes: {notes}"
                
            slide_texts.append(slide_content)
            
        structured_text = "\n".join(slide_texts)

        from ai_providers.router import router
        prompt = f"""
        Extract key financial or traction metrics from the following pitch deck text. 
        Focus on ARR, Customer Count, MRR, Gross Margin, Burn Rate.
        Return a JSON list of objects with keys: metric_name, metric_value, period, confidence (High/Medium/Low).
        
        Deck Text:
        {structured_text[:6000]} # Limit to 6k chars for prompt safety
        """
        
        result = router.execute_task("data_room_deck", prompt)
        extracted = []
        
        if result and "data" in result:
            data = result["data"]
            if isinstance(data, list):
                for item in data:
                    try:
                        extracted.append(ExtractedMetric(
                            metric_name=item.get("metric_name", "Unknown"),
                            metric_value=str(item.get("metric_value", "Unknown")),
                            source_document=document_name,
                            period=item.get("period", "Unknown"),
                            confidence=item.get("confidence", "Medium"),
                            verification_status="Extracted from Deck (Unverified)"
                        ))
                    except Exception as e:
                        logger.error(f"Error parsing deck metric {item}: {e}")
            
        return extracted
    except Exception as e:
        logger.error(f"Failed to parse pitch deck {file_path}: {e}")
        return []
